import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Helper Functions ---
    def add_title_slide(title_text, subtitle_text, footer_text=""):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text
        
        if footer_text:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = footer_text
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = content_text_list[0]
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # --- SLIDES (Team & Narrative) ---

    # Slide 1: Title
    add_title_slide(
        "KisanAI (FarmX)", 
        "Empowering Farmers with Cyber-Physical Systems",
        "Government Institute of Electronics\nDepartment: Cyber Physical System and Security"
    )

    # Slide 2: Meet the Team
    add_content_slide("Meet the Team", [
        "Project Contributors:",
        "1. B. Sreehitha",
        "2. N. Ajay Hemanth",
        "3. D. Hansika",
        "4. A. Ruchitha",
        "5. S. Bhagya Sree",
        "6. Ch. Sugreshwar"
    ])

    # Slide 3: Problem Statement
    add_content_slide("Problem Statement", [
        "Agriculture faces critical modern challenges:",
        "1. Late Disease Identification: Causes 20-30% yield loss globally.",
        "2. Knowledge Gap: Farmers lack access to scientific advice.",
        "3. Climate Uncertainty: Unpredictable weather disrupts planning.",
        "4. Soil Degradation: Improper fertilizer usage harms long-term fertility."
    ])

    # Slide 4: Our Vision
    add_content_slide("Our Vision", [
        "”To bridge the gap between advanced technology and the Indian farmer.”",
        "",
        "We envision a future where:",
        "- Every farmer has an expert agronomist in their pocket.",
        "- Decisions are data-driven, not just intuition-based.",
        "- Technology ensures sustainability and profitability."
    ])

    # Slide 5: Competitive Landscape (Comparison)
    add_content_slide("Why KisanAI? (Comparison)", [
        "Traditional Methods / Generic Apps:",
        "- Offer static, generalized advice.",
        "- Require manual data entry.",
        "- No direct disease diagnosis capability.",
        "",
        "KisanAI Advantage:",
        "- Hyper-Localized: Considers specific soil & weather.",
        "- AI-Driven: Uses Computer Vision for diagnosis.",
        "- Generative: Creates custom advice using LLMs (Gemini).",
        "- Holistic: One platform for Disease, Weather, and Market."
    ])

    # Slide 6: The Solution - KisanAI Ecosystem
    add_content_slide("The Solution: KisanAI Ecosystem", [
        "A comprehensive platform integrating multiple technologies:",
        "1. Visual Sening: For disease and soil analysis.",
        "2. GenAI Intelligence: For reasoning and advisory.",
        "3. Real-time Data: Weather and Market APIs.",
        "4. Accessible Interface: Designed for simplicity."
    ])

    # Slide 7: Core Feature: AI Disease Detection
    add_content_slide("Core Feature: AI Disease Detection", [
        "Problem: Visual diagnosis is difficult and error-prone.",
        "Solution: Deep Learning (CNN) powered image analysis.",
        "Tech: PyTorch Model (EfficientNet/ResNet).",
        "Workflow: Snap Photo -> Upload -> Instant Diagnosis -> Cure.",
        "Result: Early detection prevents spread."
    ])

    # Slide 8: Core Feature: Smart Advisory (Gemini)
    add_content_slide("Core Feature: Smart Advisory System", [
        "Powered by Google Gemini 2.5 Flash.",
        "Not a chatbot, but a 'Reasoning Engine'.",
        "Context-Aware: Analyzes Weather + Soil + Disease inputs together.",
        "Output: Actionable plans for Watering, Fertilizers, and Care.",
        "Safety: Includes fallback logic for reliability."
    ])

    # Slide 9: Core Feature: Soil & Fertilizer Analysis
    add_content_slide("Core Feature: Soil & Fertilizer Management", [
        "Soil Classification: Automated via image analysis.",
        "Nutrient Management: Recommends NPK based on soil type.",
        "Sustainability: Prevents over-fertilization.",
        "Specifics: Tailors advice for crops (e.g., Nitrogen for Leafy crops)."
    ])

    # Slide 10: Core Feature: Weather Engine
    add_content_slide("Core Feature: Weather Integration", [
        "Real-time hyper-local weather tracking.",
        "Crucial for day-to-day farming decisions.",
        "Integration: Directly influences the AI Advisory output.",
        "Example: 'Rain predicted? -> Skip watering today.'"
    ])

    # Slide 11: Market & Community Support
    add_content_slide("Market & Community Support", [
        "Economic Security: Access to Mandi prices (Market).",
        "Financial Aid: Information on Government Schemes.",
        "Crop Protection: Database of safe pesticides and biological controls.",
        "Goal: Protecting the farmer's wallet as well as their crop."
    ])

    # Slide 12: Technology Stack
    add_content_slide("Technology Stack", [
        "Frontend: HTML5, Vanilla CSS, JavaScript.",
        "Backend: FastAPI (Python) - High performance Async.",
        "AI/ML: PyTorch (Vision) + Google GenAI (LLM).",
        "Database: SQLite (Efficient structured storage).",
        "Deployment: Docker / Cloud-ready."
    ])

    # Slide 13: System Architecture
    add_content_slide("System Architecture", [
        "[User] <-> [Frontend UI]",
        "       |REST API",
        "       v",
        "[FastAPI Backend]",
        "   /      |       \\",
        "[PyTorch] [SQLite] [Gemini API]",
        "   (Vision)  (Data)   (Reasoning)"
    ])

    # Slide 14: Social & Economic Impact
    add_content_slide("Social & Economic Impact", [
        "Social: Democratizing access to expert knowledge.",
        "Economic: Reducing crop loss = Higher income.",
        "Environmental: Precise resource usage (Water/Fertilizer).",
        "Educational: Teaching best practices through technology."
    ])

    # Slide 15: Future Roadmap
    add_content_slide("Future Roadmap", [
        "1. Vernacular Voice Support: Talking to the app in local dialects.",
        "2. Offline Mode: AI inference on the edge (mobile device).",
        "3. IoT Integration: Connecting with smart moisture sensors.",
        "4. Drone Integration: Aerial field analysis."
    ])

    # Slide 16: Conclusion
    add_content_slide("Conclusion", [
        "KisanAI is a step towards Cyber-Physical Systems in Agriculture.",
        "We combine Sensing, Computing, and Control (Advisory).",
        "Produced by the students of Government Institute of Electronics.",
        "Thank You!"
    ])

    # Save
    output_path = "KisanAI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
